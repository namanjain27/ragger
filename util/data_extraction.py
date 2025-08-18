# input file
# use appropriate function
# return extracted data

import fitz  # PyMuPDF (install via pip install PyMuPDF)
import docx  # python-docx (install via pip install python-docx)
from pptx import Presentation  # python-pptx (pip install python-pptx)
import os
import io
from PIL import Image
import base64
from typing import List, Tuple, Dict, Any


# OCR availability check
try:
    import pytesseract
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False
    print("Warning: pytesseract not available. OCR functionality will be limited.")

try:
    from google.cloud import vision
    GOOGLE_VISION_AVAILABLE = True
except ImportError:
    GOOGLE_VISION_AVAILABLE = False
    print("Warning: Google Vision API not available. Using Tesseract for OCR.")

# OCR impl
def extract_text_from_image(image_bytes: bytes, use_google_vision: bool = True) -> str:
    """Extract text from image bytes using OCR"""
    try:
        if use_google_vision and GOOGLE_VISION_AVAILABLE:
            return _extract_text_google_vision(image_bytes)
        elif TESSERACT_AVAILABLE:
            return _extract_text_tesseract(image_bytes)
        else:
            return "[OCR not available - no text extracted from image]"
    except Exception as e:
        print(f"OCR failed: {e}")
        return "[OCR failed - no text extracted from image]"

def _extract_text_google_vision(image_bytes: bytes) -> str:
    """Extract text using Google Vision API"""
    client = vision.ImageAnnotatorClient()
    image = vision.Image(content=image_bytes)
    response = client.text_detection(image=image)
    texts = response.text_annotations
    
    if texts:
        return texts[0].description
    return ""

def _extract_text_tesseract(image_bytes: bytes) -> str:
    """Extract text using Tesseract OCR"""
    image = Image.open(io.BytesIO(image_bytes))
    return pytesseract.image_to_string(image)

# validation function for allowed file types
def validate_file(file_path: str, max_size_mb: int = 50) -> Dict[str, Any]:
    """Validate file type and size"""
    if not os.path.exists(file_path):
        return {"valid": False, "error": "File does not exist"}
    
    file_size = os.path.getsize(file_path) / (1024 * 1024)  # Size in MB
    if file_size > max_size_mb:
        return {"valid": False, "error": f"File size ({file_size:.1f}MB) exceeds limit ({max_size_mb}MB)"}
    
    file_ext = os.path.splitext(file_path)[1].lower()
    supported_exts = ['.pdf', '.docx', '.pptx', '.txt', '.jpg', '.jpeg', '.png', '.bmp', '.tiff']
    
    if file_ext not in supported_exts:
        return {"valid": False, "error": f"Unsupported file type: {file_ext}"}
    
    return {"valid": True, "size_mb": file_size, "type": file_ext}

# extraction functions for specific file types
def extract_from_pdf(pdf_path: str) -> Tuple[str, List[Tuple[str, str]], List[str]]:
    """Extract text, images with OCR, and links from PDF"""
    validation = validate_file(pdf_path)
    if not validation["valid"]:
        raise ValueError(f"Invalid PDF file: {validation['error']}")
    
    doc = fitz.open(pdf_path)
    all_text = ""
    images = []
    links = []
    
    for i, page in enumerate(doc):
        # Extract regular text
        page_text = page.get_text()
        all_text += f"\n--- Page {i+1} ---\n{page_text}"
        
        # Extract images and perform OCR
        image_list = page.get_images(full=True)
        for img_index, img in enumerate(image_list):
            try:
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                image_name = f"{os.path.splitext(os.path.basename(pdf_path))[0]}_page{i+1}_img{img_index+1}.{image_ext}"
                
                # Perform OCR on the image
                ocr_text = extract_text_from_image(image_bytes)
                if ocr_text.strip():
                    all_text += f"\n[Image OCR - {image_name}]: {ocr_text}"
                
                images.append((image_name, ocr_text))
            except Exception as e:
                print(f"Error processing image {img_index} on page {i+1}: {e}")
        
        # Extract links
        links.extend([link.get('uri', '') for link in page.get_links() if link.get('uri')])
    
    doc.close()
    return all_text, images, links

def extract_from_docx(docx_path: str) -> Tuple[str, List[Tuple[str, str]], List[str]]:
    """Extract text, images with OCR, and links from DOCX"""
    validation = validate_file(docx_path)
    if not validation["valid"]:
        raise ValueError(f"Invalid DOCX file: {validation['error']}")
    
    doc = docx.Document(docx_path)
    all_text = ""
    images = []
    links = []
    
    # Extract text from paragraphs
    for i, para in enumerate(doc.paragraphs):
        if para.text.strip():
            all_text += f"{para.text}\n"
        
        # Extract hyperlinks
        for run in para.runs:
            if hasattr(run, '_element') and run._element.hyperlink:
                hyperlink = run._element.hyperlink
                if hasattr(hyperlink, 'target_ref') and hyperlink.target_ref:
                    links.append(hyperlink.target_ref)
    
    # Extract images and perform OCR
    img_count = 0
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            try:
                image_bytes = rel.target_part.blob
                image_name = f"{os.path.splitext(os.path.basename(docx_path))[0]}_img{img_count+1}"
                
                # Perform OCR on the image
                ocr_text = extract_text_from_image(image_bytes)
                if ocr_text.strip():
                    all_text += f"\n[Image OCR - {image_name}]: {ocr_text}\n"
                
                images.append((image_name, ocr_text))
                img_count += 1
            except Exception as e:
                print(f"Error processing image {img_count}: {e}")
    
    return all_text, images, links

def extract_from_pptx(pptx_path: str) -> Tuple[str, List[Tuple[str, str]], List[str]]:
    """Extract text, images with OCR, and links from PPTX"""
    validation = validate_file(pptx_path)
    if not validation["valid"]:
        raise ValueError(f"Invalid PPTX file: {validation['error']}")
    
    prs = Presentation(pptx_path)
    all_text = ""
    images = []
    links = []
    
    for i, slide in enumerate(prs.slides):
        slide_text = f"\n--- Slide {i+1} ---\n"
        
        for shape in slide.shapes:
            # Extract text from shapes
            if hasattr(shape, "text") and shape.text.strip():
                slide_text += shape.text + "\n"
            
            # Extract images and perform OCR
            if hasattr(shape, "image") and shape.image:
                try:
                    image_bytes = shape.image.blob
                    image_name = f"{os.path.splitext(os.path.basename(pptx_path))[0]}_slide{i+1}_img"
                    
                    # Perform OCR on the image
                    ocr_text = extract_text_from_image(image_bytes)
                    if ocr_text.strip():
                        slide_text += f"\n[Image OCR - {image_name}]: {ocr_text}\n"
                    
                    images.append((image_name, ocr_text))
                except Exception as e:
                    print(f"Error processing image on slide {i+1}: {e}")
            
            # Extract hyperlinks
            if hasattr(shape, "hyperlink") and shape.hyperlink and shape.hyperlink.address:
                links.append(shape.hyperlink.address)
        
        all_text += slide_text
    
    return all_text, images, links

def extract_from_txt(txt_path: str) -> Tuple[str, List[Tuple[str, str]], List[str]]:
    """Extract text from TXT file"""
    validation = validate_file(txt_path)
    if not validation["valid"]:
        raise ValueError(f"Invalid TXT file: {validation['error']}")
    
    try:
        with open(txt_path, "r", encoding="utf-8") as f:
            text = f.read()
    except UnicodeDecodeError:
        # Try with different encoding if UTF-8 fails
        with open(txt_path, "r", encoding="latin-1") as f:
            text = f.read()
    
    # TXT files have no images/links
    return text, [], []

def extract_from_image(image_path: str) -> Tuple[str, List[Tuple[str, str]], List[str]]:
    """Extract text from standalone image using OCR"""
    validation = validate_file(image_path)
    if not validation["valid"]:
        raise ValueError(f"Invalid image file: {validation['error']}")
    
    try:
        with open(image_path, "rb") as f:
            image_bytes = f.read()
        
        image_name = os.path.splitext(os.path.basename(image_path))[0]
        ocr_text = extract_text_from_image(image_bytes)
        
        return ocr_text, [(image_name, ocr_text)], []
    except Exception as e:
        raise ValueError(f"Error processing image {image_path}: {e}")

# Primary function used for data extraction that calls secondary functions
def extract_content(file_path: str) -> Tuple[str, List[Tuple[str, str]], List[str]]:
    """Universal content extractor that determines file type and uses appropriate extraction method"""
    validation = validate_file(file_path)
    if not validation["valid"]:
        raise ValueError(f"Invalid file: {validation['error']}")
    
    file_ext = validation["type"]
    
    try:
        if file_ext == '.pdf':
            return extract_from_pdf(file_path)
        elif file_ext == '.docx':
            return extract_from_docx(file_path)
        elif file_ext == '.pptx':
            return extract_from_pptx(file_path)
        elif file_ext == '.txt':
            return extract_from_txt(file_path)
        elif file_ext in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff']:
            return extract_from_image(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_ext}")
    except Exception as e:
        raise ValueError(f"Error extracting content from {file_path}: {e}")

# Usage Examples:
# text, images, links = extract_content("input.pdf")  # Universal extractor
# pdf_text, pdf_images, pdf_links = extract_from_pdf("input.pdf")
# docx_text, docx_images, docx_links = extract_from_docx("input.docx")
# pptx_text, pptx_images, pptx_links = extract_from_pptx("input.pptx")
# txt_text, _, _ = extract_from_txt("input.txt")
# img_text, _, _ = extract_from_image("input.jpg")
