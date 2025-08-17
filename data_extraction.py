# input file
# use appropriate function
# return extracted data

import fitz  # PyMuPDF (install via pip install PyMuPDF)
import docx  # python-docx (install via pip install python-docx)
from pptx import Presentation  # python-pptx (pip install python-pptx)
import os

def extract_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    all_text = ""
    images = []
    links = []
    for i, page in enumerate(doc):
        all_text += page.get_text()
        image_list = page.get_images(full=True)
        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            image_name = f"{os.path.splitext(pdf_path)}_page{i+1}_img{img_index+1}.{image_ext}"
            images.append((image_name, image_bytes))
        links.extend(page.get_links())
    return all_text, images, links

def extract_from_docx(docx_path):
    doc = docx.Document(docx_path)
    all_text = " ".join([para.text for para in doc.paragraphs])
    images = []
    links = []
    for para in doc.paragraphs:
        for run in para.runs:
            if run.hyperlink:
                links.append(run.hyperlink.target)
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            images.append(rel.target_part.blob)
    return all_text, images, links

def extract_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    all_text = ""
    images = []
    links = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text += shape.text + "\n"
            if "image" in shape.name.lower():
                if hasattr(shape, "image") and shape.image:
                    image_bytes = shape.image.blob
                    image_ext = shape.image.ext
                    images.append((f"{pptx_path}_slide{i+1}.{image_ext}", image_bytes))
            if hasattr(shape, "hyperlink") and shape.hyperlink.address:
                links.append(shape.hyperlink.address)
    return all_text, images, links

def extract_from_txt(txt_path):
    with open(txt_path, "r", encoding="utf-8") as f:
        text = f.read()
    # TXT files have no images/links
    return text, [], []

# Usage Example:
# pdf_text, pdf_images, pdf_links = extract_from_pdf("input.pdf")
# docx_text, docx_images, docx_links = extract_from_docx("input.docx")
# pptx_text, pptx_images, pptx_links = extract_from_pptx("input.pptx")
# txt_text, _, _ = extract_from_txt("input.txt")
